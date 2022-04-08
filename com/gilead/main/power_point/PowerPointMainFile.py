'''
Created on Oct 3, 2018

@author: bkotamahanti
'''

from definitions import POWERPOINT_PROCESS, POWERPOINT_IMAGE_FILE

from pptx.util import Inches
from selenium import webdriver
from pptx.api import Presentation

class PowerPointMainClass():
    __process = POWERPOINT_PROCESS
    __image = POWERPOINT_IMAGE_FILE
    
    def __init__(self, utilsRef, loggerRef):
        self.utils = utilsRef
        self.logger = loggerRef
        self.driver = ""
        self.utils.killProcess(self.__process)
    
    def killProcess(self):
        self.utils.killProcess(self.__process)
    
    def startPowerPoint_via_StartMenu(self):
        self.utils.startProgramByName_via_StartMenu("PowerPoint")
    
    def isPowerPointProcessExist(self):
        return self.utils.isProcessExist(self.__process)
    
    def clickOpenOtherPresentationsButton(self):
        parent_element = self.driver.find_element_by_name("PowerPoint")
        backstage_element = parent_element.find_element_by_name("Backstage view")
        backstage_element.find_element_by_name("Open Other Presentations").click()
    
    def clickOpenWindowBrowseButton(self):
        parent_element = self.driver.find_element_by_name("PowerPoint")
        backstage_element = parent_element.find_element_by_name("Backstage view")
        backstage_element.find_element_by_name("Browse").click()

    def enterFileNameInOpenBrowseWindow(self, fileName):
        browse_window_element = self.driver.find_element_by_class_name("#32770")
        browse_window_element.find_element_by_class_name("Edit").click()
        browse_window_element.find_element_by_class_name("Edit").send_keys(fileName) 
    
    
    def clickOpenButtonOfOpenBrowseWindow(self):
        browse_window_element = self.driver.find_element_by_class_name("#32770")
        elements =  browse_window_element.find_elements_by_class_name("Button")
        for element in elements:
            if element.get_attribute("Name") == "Open":
                element.click()
    
    
    def clickCancelOnPowerpointOpenDialog(self):
        browse_window_element = self.driver.find_element_by_class_name("#32770")
        elements =  browse_window_element.find_elements_by_class_name("Button")
        for element in elements:
            if element.get_attribute("Name") == "Cancel":
                element.click()
    
    def clickPowerpointAppBackButton(self):
        parent_element = self.driver.find_element_by_class_name("PPTFrameClass") 
        parent_element.find_element_by_name("Back").click()
    
   
    def getPowerPointPopUpWindowElement(self):
        parent_element = self.driver.find_element_by_name("PowerPoint")
        open_dialog_element = parent_element.find_element_by_class_name("#32770")
        return open_dialog_element
       
    def enterTextInTitleTextBox(self, text):
        parent_element = self.driver.find_element_by_name("Presentation1 - PowerPoint")
        '''keeping the commented code as it is related to Office 2016 PPT. For reference'''
#         workspace_element = parent_element.find_element_by_name("PowerPoint Edit View - [Presentation1]")
        workspace_element = parent_element.find_element_by_class_name("MDIClient")
        slide_pane = workspace_element.find_element_by_name("").send_keys(text)
#         slide_pane.find_element_by_name("Slide 1 - ").find_element_by_name("Title TextBox").send_keys(text)
    
    def enterTextInTextBox(self, text):
        parent_element = self.driver.find_element_by_name("Presentation1 - PowerPoint")
        workspace_element = parent_element.find_element_by_name("PowerPoint Edit View - [Presentation1]")
        slide_pane = workspace_element.find_element_by_name("Slide")
        slide_pane.find_element_by_name("Slide 1 - Hello World!!").find_element_by_name("TextBox").send_keys(text)
    
    def connectPowerPointToWebDriver(self):
        self.driver = webdriver.Remote(command_executor="http://localhost:9999",
                                       desired_capabilities={
#                                             'app': r'C:\Program Files (x86)\Microsoft Office\root\Office16\POWERPNT.EXE',
                                            'app': r'POWERPNT.EXE',
                                           'debugConnectToRunningApp': 'true',
                                           "args": "-port 345"
                                        })
    
    def openPowerpoint(self, file):
        self.utils.openApplication(file)
        
        
    def createPowerPointPresentation(self, fileName):
        self.prs = Presentation()
        bullet_slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = 'Adding a Bullet Slide'
        tf = body_shape.text_frame
        tf.text = 'Find the bullet slide layout'
        p = tf.add_paragraph()
        p.text = 'Use _TextFrame.text for first bullet'
        p.level = 1
        p = tf.add_paragraph()
        p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
        p.level = 2
        left = top = Inches(4)
        slide.shapes.add_picture(self.__image, left, top)
        self.prs.save(fileName)
        self.logger.info("Created power point at {}".format(fileName))
    
    def getPowerPointMainWindowSlideElement(self, nameOfMainWindow):
        name = nameOfMainWindow+" - PowerPoint"
        parent_element = self.driver.find_element_by_name(name)
        MsoDockTop_element = parent_element.find_element_by_name("Slide")
        return MsoDockTop_element
    
    def clickBlankPresentation(self):
        parent_element = self.driver.find_element_by_name("PowerPoint")
        parent_element.find_element_by_name("Blank Presentation").click()
    
    def clickDontSaveInMsPopUp(self):
        parent_element = self.driver.find_element_by_name("Presentation1 - PowerPoint")
        pane_element = parent_element.find_element_by_name("Microsoft PowerPoint")
        pane_element.find_element_by_name("Don't Save").click()
        
    def clickCloseButtonOfPowerPointApp(self):
        parent_element = self.driver.find_element_by_class_name("PPTFrameClass")
        doctop_element = parent_element.find_element_by_name("MsoDockTop")
        doctop_element.find_element_by_name("Close").click()
    
    def clickCloseButtonOfPowerPointAppWhenNoPresentation(self):
        parent_element = self.driver.find_element_by_class_name("PPTFrameClass")
        parent_element.find_element_by_name("Close").click()
    
    
    def clickOKOnPowepointOpenDialog(self):
        open_dialog_element = self.getPowerPointPopUpWindowElement()
        open_dialog_element.find_element_by_name("OK").click()
    
    
        
    def isPowerPointElementExist(self, webelement, **property1):
        return self.utils.isElementPresent( webelement,  **property1)
    
    def isPowerPointElementTextExist(self, webelement, text , **property1):
        return self.utils.isElementTextPresent(webelement, text , **property1)